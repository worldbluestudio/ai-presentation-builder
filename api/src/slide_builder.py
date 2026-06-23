"""
スライドビルダーモジュール
6種類のスライドレイアウトをプロ品質で構築する。
テキストの可視性を最優先に設計。
"""

import os
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

from src.theme import Theme, hex_to_rgb
from src.decorations import (
    SLIDE_WIDTH, SLIDE_HEIGHT,
    add_solid_background, add_gradient_background,
    add_accent_bar_top, add_accent_bar_left,
    add_separator_line, add_overlay_rectangle,
    add_image_with_shadow, add_page_number,
    add_decorative_circle, add_bullet_icon,
)

# lxml namespace
_NSMAP = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}


# =========================================================================
# テキスト可視性ヘルパー
# =========================================================================

def _set_shape_opacity(shape, opacity: float):
    """シェイプの塗りつぶしに不透明度を設定する。"""
    opacity_val = int(opacity * 100000)
    solid_fill = shape._element.find('.//a:solidFill', _NSMAP)
    if solid_fill is not None:
        srgb = solid_fill.find('a:srgbClr', _NSMAP)
        if srgb is not None:
            # 既存のalphaを削除
            for old_alpha in srgb.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}alpha'):
                srgb.remove(old_alpha)
            alpha = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha.set('val', str(opacity_val))


def _add_text_backdrop(slide, left, top, width, height,
                       color_rgb: RGBColor, opacity: float = 0.85,
                       corner_radius: bool = True):
    """
    テキストの背後に半透明の背板を配置する。
    テキストの可視性を確実に保証するための核心機能。
    """
    if corner_radius:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height,
        )
        shape.adjustments[0] = 0.04
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height,
        )

    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb
    _set_shape_opacity(shape, opacity)
    return shape


def _get_contrast_text_color(theme: Theme, on_image: bool = False) -> RGBColor:
    """
    背景に対してコントラストが確保されたテキスト色を返す。
    画像上のテキストは常に白を使用する。
    """
    if on_image:
        return RGBColor(0xFF, 0xFF, 0xFF)
    return theme.rgb("title_color")


def _get_contrast_sub_color(theme: Theme, on_image: bool = False) -> RGBColor:
    """画像上のサブテキスト用のコントラスト色。"""
    if on_image:
        return RGBColor(0xE0, 0xE0, 0xE0)
    return theme.rgb("text_light")


# =========================================================================
# テキスト配置ヘルパー
# =========================================================================

def _add_text_box(slide, left, top, width, height, text: str,
                  font_size, font_color_rgb, bold=False,
                  alignment=PP_ALIGN.LEFT, font_name="游ゴシック",
                  word_wrap=True, vertical_anchor=MSO_ANCHOR.TOP):
    """テキストボックスを追加するヘルパー。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color_rgb
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment

    # 縦方向のアンカー
    body_pr = txBox._element.find('.//a:bodyPr', _NSMAP)
    if body_pr is not None:
        anchor_map = {
            MSO_ANCHOR.TOP: 't',
            MSO_ANCHOR.MIDDLE: 'ctr',
            MSO_ANCHOR.BOTTOM: 'b',
        }
        body_pr.set('anchor', anchor_map.get(vertical_anchor, 't'))

    return txBox


def _add_multiline_text(slide, left, top, width, height, lines: list[str],
                        font_size, font_color_rgb, line_spacing: float = 1.5,
                        font_name="游ゴシック", theme: Theme = None):
    """複数行のテキストボックスを追加する。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line
        p.font.size = font_size
        p.font.color.rgb = font_color_rgb
        p.font.name = font_name
        p.space_after = Pt(8)
        p.line_spacing = line_spacing

    return txBox


def _add_styled_bullets(slide, left, top, width, height,
                        bullets: list[dict], theme: Theme):
    """
    スタイル付き箇条書きを追加する。
    bullets: [{"title": "見出し", "desc": "説明文"}, ...] or ["テキスト", ...]
    """
    current_top = top
    item_spacing = Pt(12)

    for item in bullets:
        if isinstance(item, dict):
            title_text = item.get("title", "")
            desc_text = item.get("desc", "")
        else:
            title_text = str(item)
            desc_text = ""

        # アイコンドット
        add_bullet_icon(slide, theme, left, current_top + Pt(5))

        # タイトル行
        icon_offset = Pt(20)
        _add_text_box(
            slide, left + icon_offset, current_top,
            width - icon_offset, Pt(24),
            title_text,
            font_size=theme.fonts.body_size,
            font_color_rgb=theme.rgb("title_color"),
            bold=True,
            font_name=theme.fonts.body_ja,
        )
        current_top += Pt(26)

        # 説明文（あれば）
        if desc_text:
            _add_text_box(
                slide, left + icon_offset, current_top,
                width - icon_offset, Pt(20),
                desc_text,
                font_size=Pt(13),
                font_color_rgb=theme.rgb("text_color"),
                font_name=theme.fonts.body_ja,
            )
            current_top += Pt(24)

        current_top += item_spacing


# =========================================================================
# スライドレイアウト関数
# =========================================================================

def build_title_slide(prs, slide_data: dict, theme: Theme):
    """
    タイトルスライド: 背景画像 + 強力なオーバーレイ + バックドロップ付きタイトル
    """
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # 背景画像（あれば）
    image_path = slide_data.get("image")
    has_image = image_path and os.path.exists(image_path)

    if has_image:
        slide.shapes.add_picture(
            image_path, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT,
        )
        # 全面オーバーレイ（強めに暗くする）
        overlay = add_overlay_rectangle(
            slide, theme,
            Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT,
        )
        # オーバーレイを確実に濃くする
        _set_shape_opacity(overlay, 0.80)
    else:
        add_gradient_background(slide, theme)

    # 装飾サークル（左上と右下）
    add_decorative_circle(slide, theme, Inches(-1.5), Inches(-1.5), Inches(5), opacity=0.08)
    add_decorative_circle(slide, theme, SLIDE_WIDTH - Inches(3), SLIDE_HEIGHT - Inches(3), Inches(6), opacity=0.06)

    # テキスト用バックドロップ（タイトル+サブタイトル領域をカバー）
    title_top = Inches(2.0)
    backdrop_color = theme.rgb("bg_primary") if not has_image else theme.rgb("overlay")
    _add_text_backdrop(
        slide,
        Inches(1.0), title_top - Inches(0.2),
        Inches(10.5), Inches(3.2),
        color_rgb=backdrop_color,
        opacity=0.88 if has_image else 0.70,
        corner_radius=True,
    )

    # テキスト色の決定（画像上では白）
    title_color = _get_contrast_text_color(theme, on_image=has_image)
    sub_color = _get_contrast_sub_color(theme, on_image=has_image)

    # タイトル
    title = slide_data.get("title", "")
    _add_text_box(
        slide, Inches(1.5), title_top,
        Inches(10), Inches(1.5),
        title,
        font_size=Pt(44),
        font_color_rgb=title_color,
        bold=True,
        alignment=PP_ALIGN.LEFT,
        font_name=theme.fonts.title_ja,
    )

    # セパレーターライン
    add_separator_line(
        slide, theme,
        Inches(1.5), title_top + Inches(1.6),
        Inches(3),
    )

    # サブタイトル
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        _add_text_box(
            slide, Inches(1.5), title_top + Inches(1.9),
            Inches(10), Inches(0.8),
            subtitle,
            font_size=theme.fonts.subtitle_size,
            font_color_rgb=sub_color,
            alignment=PP_ALIGN.LEFT,
            font_name=theme.fonts.body_ja,
        )

    # 下部にアクセントバー
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), SLIDE_HEIGHT - Pt(6),
        SLIDE_WIDTH, Pt(6),
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.rgb("accent")

    return slide


def build_content_slide(prs, slide_data: dict, theme: Theme, page_num: int = 0):
    """
    コンテンツスライド: 左にテキスト + 右に画像
    タイトルにバックドロップを配置して可視性を確保。
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    add_gradient_background(slide, theme)

    # 上部アクセントバー
    add_accent_bar_top(slide, theme, height_pt=6)

    # 左端アクセント（細いライン）
    add_accent_bar_left(slide, theme, width_pt=4)

    # タイトル用バックドロップ
    _add_text_backdrop(
        slide,
        Inches(0.5), Inches(0.2),
        Inches(12), Inches(0.9),
        color_rgb=theme.rgb("bg_secondary"),
        opacity=0.90,
        corner_radius=True,
    )

    # タイトル
    title = slide_data.get("title", "")
    _add_text_box(
        slide, Inches(0.8), Inches(0.4),
        Inches(8), Inches(0.8),
        title,
        font_size=theme.fonts.heading_size,
        font_color_rgb=theme.rgb("title_color"),
        bold=True,
        font_name=theme.fonts.title_ja,
    )

    # タイトル下セパレーター
    add_separator_line(
        slide, theme,
        Inches(0.8), Inches(1.15),
        Inches(2.5),
    )

    # 箇条書き / コンテンツエリア（左側）
    bullets = slide_data.get("bullets", [])
    if bullets:
        _add_styled_bullets(
            slide,
            left=Inches(1.0),
            top=Inches(1.6),
            width=Inches(5.5),
            height=Inches(5),
            bullets=bullets,
            theme=theme,
        )

    # 画像（右側）
    image_path = slide_data.get("image")
    if image_path and os.path.exists(image_path):
        img_left = Inches(7.5)
        img_top = Inches(1.5)
        img_width = Inches(5.2)
        img_height = Inches(5.2)
        add_image_with_shadow(slide, image_path, img_left, img_top, img_width, img_height)

    # ページ番号
    if page_num > 0:
        add_page_number(slide, theme, page_num)

    return slide


def build_image_full_slide(prs, slide_data: dict, theme: Theme, page_num: int = 0):
    """
    フルブリード画像スライド: 画像全面 + 強力なキャプション帯
    テキスト領域は必ず高コントラストの背板で保護する。
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景（画像がない場合のフォールバック）
    add_gradient_background(slide, theme)

    # フルブリード画像
    image_path = slide_data.get("image")
    has_image = image_path and os.path.exists(image_path)
    if has_image:
        slide.shapes.add_picture(
            image_path, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT,
        )

    # 下部グラデーション帯（強めのオーバーレイ）
    caption_height = Inches(2.5)
    overlay = add_overlay_rectangle(
        slide, theme,
        Emu(0), SLIDE_HEIGHT - caption_height,
        SLIDE_WIDTH, caption_height,
    )
    _set_shape_opacity(overlay, 0.90)

    # テキスト色（画像上では白）
    title_color = _get_contrast_text_color(theme, on_image=has_image)
    sub_color = _get_contrast_sub_color(theme, on_image=has_image)

    # キャプションタイトル
    title = slide_data.get("title", "")
    if title:
        _add_text_box(
            slide, Inches(1.2), SLIDE_HEIGHT - Inches(2.2),
            Inches(10), Inches(0.8),
            title,
            font_size=Pt(30),
            font_color_rgb=title_color,
            bold=True,
            font_name=theme.fonts.title_ja,
        )

    # キャプション説明文
    caption = slide_data.get("caption", "")
    if caption:
        _add_text_box(
            slide, Inches(1.2), SLIDE_HEIGHT - Inches(1.2),
            Inches(10), Inches(0.6),
            caption,
            font_size=Pt(16),
            font_color_rgb=sub_color,
            font_name=theme.fonts.body_ja,
        )

    # ページ番号
    if page_num > 0:
        add_page_number(slide, theme, page_num)

    return slide


def build_two_column_slide(prs, slide_data: dict, theme: Theme, page_num: int = 0):
    """
    2カラムスライド: 左右に情報を分割
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    add_gradient_background(slide, theme)

    # 上部アクセントバー
    add_accent_bar_top(slide, theme, height_pt=6)

    # タイトル用バックドロップ
    _add_text_backdrop(
        slide,
        Inches(0.5), Inches(0.2),
        Inches(12), Inches(0.9),
        color_rgb=theme.rgb("bg_secondary"),
        opacity=0.90,
        corner_radius=True,
    )

    # タイトル
    title = slide_data.get("title", "")
    _add_text_box(
        slide, Inches(0.8), Inches(0.4),
        Inches(11), Inches(0.8),
        title,
        font_size=theme.fonts.heading_size,
        font_color_rgb=theme.rgb("title_color"),
        bold=True,
        font_name=theme.fonts.title_ja,
    )

    # セパレーター
    add_separator_line(slide, theme, Inches(0.8), Inches(1.15), Inches(2.5))

    # 左カラム
    left_content = slide_data.get("left", {})
    left_text = left_content.get("text", [])
    if isinstance(left_text, str):
        left_text = [left_text]

    if left_text:
        _add_multiline_text(
            slide, Inches(0.8), Inches(1.6),
            Inches(5.5), Inches(5),
            left_text,
            font_size=theme.fonts.body_size,
            font_color_rgb=theme.rgb("text_color"),
            font_name=theme.fonts.body_ja,
            theme=theme,
        )

    left_image = left_content.get("image")
    if left_image and os.path.exists(left_image):
        add_image_with_shadow(
            slide, left_image,
            Inches(0.8), Inches(1.6),
            Inches(5.5), Inches(5),
        )

    # 右カラム
    right_content = slide_data.get("right", {})
    right_text = right_content.get("text", [])
    if isinstance(right_text, str):
        right_text = [right_text]

    if right_text:
        _add_multiline_text(
            slide, Inches(7.0), Inches(1.6),
            Inches(5.5), Inches(5),
            right_text,
            font_size=theme.fonts.body_size,
            font_color_rgb=theme.rgb("text_color"),
            font_name=theme.fonts.body_ja,
            theme=theme,
        )

    right_image = right_content.get("image")
    if right_image and os.path.exists(right_image):
        add_image_with_shadow(
            slide, right_image,
            Inches(7.0), Inches(1.6),
            Inches(5.5), Inches(5),
        )

    # ページ番号
    if page_num > 0:
        add_page_number(slide, theme, page_num)

    return slide


def build_key_message_slide(prs, slide_data: dict, theme: Theme, page_num: int = 0):
    """
    キーメッセージスライド: 中央に大きなメッセージ + バックドロップ + 装飾
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    add_gradient_background(slide, theme)

    # 装飾サークル（複数配置）
    add_decorative_circle(slide, theme, Inches(0.5), Inches(0.5), Inches(3), opacity=0.07)
    add_decorative_circle(slide, theme, Inches(9), Inches(4), Inches(4), opacity=0.05)
    add_decorative_circle(slide, theme, Inches(10.5), Inches(0.8), Inches(2), opacity=0.10)

    # メッセージ用バックドロップ（中央の大きな背板）
    _add_text_backdrop(
        slide,
        Inches(1.0), Inches(1.5),
        Inches(11), Inches(4.5),
        color_rgb=theme.rgb("bg_secondary"),
        opacity=0.80,
        corner_radius=True,
    )

    # メインメッセージ
    message = slide_data.get("message", "")
    _add_text_box(
        slide, Inches(1.5), Inches(2.0),
        Inches(10), Inches(2.0),
        message,
        font_size=Pt(36),
        font_color_rgb=theme.rgb("title_color"),
        bold=True,
        alignment=PP_ALIGN.CENTER,
        font_name=theme.fonts.title_ja,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )

    # セパレーター（中央）
    add_separator_line(slide, theme, Inches(5.0), Inches(4.2), Inches(3))

    # 補足テキスト
    sub_message = slide_data.get("sub_message", "")
    if sub_message:
        _add_text_box(
            slide, Inches(2), Inches(4.6),
            Inches(9), Inches(1.2),
            sub_message,
            font_size=Pt(18),
            font_color_rgb=theme.rgb("text_color"),
            alignment=PP_ALIGN.CENTER,
            font_name=theme.fonts.body_ja,
        )

    # ページ番号
    if page_num > 0:
        add_page_number(slide, theme, page_num)

    return slide


def build_ending_slide(prs, slide_data: dict, theme: Theme):
    """
    エンディングスライド: 背景画像 + 強力なバックドロップ + メッセージ
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景画像（あれば）
    image_path = slide_data.get("image")
    has_image = image_path and os.path.exists(image_path)

    if has_image:
        slide.shapes.add_picture(
            image_path, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT,
        )
        # 全面オーバーレイ（強め）
        overlay = add_overlay_rectangle(
            slide, theme,
            Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT,
        )
        _set_shape_opacity(overlay, 0.80)
    else:
        add_gradient_background(slide, theme)

    # 装飾サークル
    add_decorative_circle(slide, theme, Inches(5.5), Inches(-0.5), Inches(3), opacity=0.08)

    # テキスト用バックドロップ（中央に大きく）
    backdrop_color = theme.rgb("bg_primary") if not has_image else theme.rgb("overlay")
    _add_text_backdrop(
        slide,
        Inches(2.0), Inches(1.5),
        Inches(9), Inches(4.5),
        color_rgb=backdrop_color,
        opacity=0.85 if has_image else 0.70,
        corner_radius=True,
    )

    # テキスト色の決定
    title_color = _get_contrast_text_color(theme, on_image=has_image)
    sub_color = _get_contrast_sub_color(theme, on_image=has_image)

    # メインメッセージ
    message = slide_data.get("message", "Thank You")
    _add_text_box(
        slide, Inches(1), Inches(2.0),
        Inches(11), Inches(1.5),
        message,
        font_size=Pt(48),
        font_color_rgb=title_color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        font_name=theme.fonts.title_ja,
    )

    # セパレーター
    add_separator_line(slide, theme, Inches(5.2), Inches(3.7), Inches(3))

    # サブメッセージ / コンタクト
    sub_message = slide_data.get("sub_message", "")
    if sub_message:
        _add_text_box(
            slide, Inches(2), Inches(4.2),
            Inches(9), Inches(1.5),
            sub_message,
            font_size=Pt(18),
            font_color_rgb=sub_color,
            alignment=PP_ALIGN.CENTER,
            font_name=theme.fonts.body_ja,
        )

    # 下部アクセントバー
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), SLIDE_HEIGHT - Pt(6),
        SLIDE_WIDTH, Pt(6),
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.rgb("accent")

    return slide


# =========================================================================
# スライドタイプ → ビルダー関数のマッピング
# =========================================================================
SLIDE_BUILDERS = {
    "title": build_title_slide,
    "content": build_content_slide,
    "image_full": build_image_full_slide,
    "two_column": build_two_column_slide,
    "key_message": build_key_message_slide,
    "ending": build_ending_slide,
}
