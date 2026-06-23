"""
装飾要素モジュール
スライドに配置するアクセントバー、セパレーター、オーバーレイ等を生成する。
"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from src.theme import Theme, hex_to_rgb


# =========================================================================
# スライドサイズ定数 (16:9)
# =========================================================================
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_solid_background(slide, theme: Theme):
    """スライド全体にソリッド背景を設定する。"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = theme.rgb("bg_primary")


def add_gradient_background(slide, theme: Theme):
    """スライド全体にグラデーション背景を設定する。"""
    # python-pptx ではスライド背景のグラデーション設定が限定的なため、
    # フルサイズのシェイプを最背面に配置してグラデーションを実現する。
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0),
        SLIDE_WIDTH, SLIDE_HEIGHT,
    )
    shape.line.fill.background()  # 枠線なし

    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = theme.rgb("bg_primary")
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = theme.rgb("bg_secondary")
    fill.gradient_stops[1].position = 1.0

    # 最背面に移動
    sp = shape._element
    sp.getparent().insert(0, sp)


def add_accent_bar_top(slide, theme: Theme, height_pt: int = 8):
    """スライド上部にアクセントカラーのバーを配置する。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0),
        SLIDE_WIDTH, Pt(height_pt),
    )
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = theme.rgb("accent")


def add_accent_bar_left(slide, theme: Theme, width_pt: int = 6):
    """スライド左端にアクセントカラーの縦バーを配置する。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0),
        Pt(width_pt), SLIDE_HEIGHT,
    )
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = theme.rgb("accent")


def add_separator_line(slide, theme: Theme, left, top, width):
    """水平のセパレーターラインを追加する。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top,
        width, Pt(3),
    )
    shape.line.fill.background()
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = theme.rgb("accent")
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = theme.rgb("accent_secondary")
    fill.gradient_stops[1].position = 1.0


def add_overlay_rectangle(slide, theme: Theme, left, top, width, height,
                          corner_radius: Emu | None = None):
    """半透明のオーバーレイ矩形を追加する。"""
    if corner_radius:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height,
        )
        shape.adjustments[0] = 0.05  # 角丸の度合い
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height,
        )

    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = theme.rgb("overlay")

    # 不透明度を設定 (0=完全透明, 100000=完全不透過)
    opacity_val = int(theme.colors.overlay_opacity * 100000)
    from lxml import etree
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    }
    solid_fill = shape._element.find('.//a:solidFill', nsmap)
    if solid_fill is not None:
        srgb = solid_fill.find('a:srgbClr', nsmap)
        if srgb is not None:
            alpha = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha.set('val', str(opacity_val))

    return shape


def add_image_with_shadow(slide, image_path: str, left, top, width, height):
    """画像を追加し、ドロップシャドウ風の装飾を付ける。"""
    # シャドウ用の矩形（画像の少し下・右にずらして配置）
    shadow_offset = Pt(4)
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + shadow_offset, top + shadow_offset,
        width, height,
    )
    shadow.adjustments[0] = 0.03
    shadow.line.fill.background()
    fill = shadow.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    # シャドウの透明度
    from lxml import etree
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    solid_fill = shadow._element.find('.//a:solidFill', nsmap)
    if solid_fill is not None:
        srgb = solid_fill.find('a:srgbClr', nsmap)
        if srgb is not None:
            alpha = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha.set('val', '20000')  # 20% 不透明

    # 画像本体
    pic = slide.shapes.add_picture(image_path, left, top, width, height)
    return pic


def add_page_number(slide, theme: Theme, number: int):
    """右下にページ番号を配置する。"""
    from pptx.util import Inches, Pt
    left = SLIDE_WIDTH - Inches(0.8)
    top = SLIDE_HEIGHT - Inches(0.5)

    txBox = slide.shapes.add_textbox(left, top, Inches(0.6), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = theme.fonts.page_number_size
    p.font.color.rgb = theme.rgb("text_light")
    p.alignment = PP_ALIGN.RIGHT


def add_decorative_circle(slide, theme: Theme, left, top, size, opacity: float = 0.15):
    """装飾用の半透明サークルを追加する。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left, top, size, size,
    )
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = theme.rgb("accent")

    opacity_val = int(opacity * 100000)
    from lxml import etree
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    solid_fill = shape._element.find('.//a:solidFill', nsmap)
    if solid_fill is not None:
        srgb = solid_fill.find('a:srgbClr', nsmap)
        if srgb is not None:
            alpha = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha.set('val', str(opacity_val))

    return shape


def add_decorative_line(slide, theme: Theme, start_x, start_y, end_x, end_y, width_pt: float = 1.5):
    """装飾用のラインを追加する。"""
    from pptx.util import Pt, Emu
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_x, start_y,
        end_x, end_y,
    )
    connector.line.color.rgb = theme.rgb("accent")
    connector.line.width = Pt(width_pt)
    return connector


def add_bullet_icon(slide, theme: Theme, left, top, size=Pt(8)):
    """箇条書きの代わりに使うアイコン風の小さな図形。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, size, size,
    )
    shape.adjustments[0] = 0.3  # かなり丸い
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = theme.rgb("accent")
    return shape
